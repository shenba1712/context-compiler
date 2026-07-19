#!/usr/bin/env python3
"""Rebuild context-compiler-pitch.pptx to match the Forest map web theme.

Tokens mirror public/style.css / docs/specs/08-design-system.md.
Fonts: Georgia ≈ Fraunces, Helvetica ≈ DM Sans, Menlo ≈ IBM Plex Mono
(system-safe pptx substitutes; install the web fonts for a closer match).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "context-compiler-pitch.pptx"

# --- Forest map tokens (public/style.css) ---
PAPER = RGBColor(0xE8, 0xEC, 0xE9)
INK = RGBColor(0x1A, 0x22, 0x1E)
FOREST = RGBColor(0x1F, 0x5C, 0x42)
FOREST_DEEP = RGBColor(0x14, 0x3D, 0x2C)
SURFACE = RGBColor(0xF3, 0xF6, 0xF4)
BORDER = RGBColor(0xC8, 0xD2, 0xCB)
BORDER2 = RGBColor(0xA8, 0xB9, 0xAF)
MUTED = RGBColor(0x4A, 0x58, 0x50)
FAINT = RGBColor(0x5C, 0x6B, 0x62)
WASTE = RGBColor(0x9A, 0xAB, 0x9F)
COMPILED = RGBColor(0x8F, 0xCE, 0xB0)
PLANE = RGBColor(0x1F, 0x5C, 0x42)
PLANE_INK = RGBColor(0xE8, 0xF0, 0xEB)
PLANE_MUTED = RGBColor(0xB8, 0xD0, 0xC4)

FONT_BRAND = "Georgia"  # Fraunces on web
FONT_BODY = "Helvetica"  # DM Sans on web
FONT_MONO = "Menlo"  # IBM Plex Mono on web

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)


def _set_run(run, *, size, bold=False, color=INK, font=FONT_BODY):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Force east-asian / complex script face too (avoids Calibri fallback on some hosts)
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.insert(0, el)
        el.set("typeface", font)


def _add_text(shape, paragraphs, *, valign=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[valign])
    except Exception:
        pass
    first = True
    for spec in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(spec.get("space_before", 0))
        p.space_after = Pt(spec.get("space_after", 0))
        run = p.add_run()
        run.text = spec["text"]
        _set_run(
            run,
            size=spec.get("size", 16),
            bold=spec.get("bold", False),
            color=spec.get("color", INK),
            font=spec.get("font", FONT_BODY),
        )


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _hairline(shape, color: RGBColor = BORDER):
    shape.line.color.rgb = color
    shape.line.width = Pt(1)


def _blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # Paper background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, PAPER)
    # Send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def _rect(slide, left, top, width, height, fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill(sh, fill)
    if line is not None:
        _hairline(sh, line)
    else:
        sh.line.fill.background()
    return sh


def _textbox(slide, left, top, width, height, paragraphs, *, valign=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(left, top, width, height)
    _add_text(sh, paragraphs, valign=valign)
    return sh


def _accent_bar(slide, left, top, width, height=Inches(0.05)):
    return _rect(slide, left, top, width, height, FOREST)


def slide_title(prs: Presentation):
    slide = _blank_slide(prs)
    # Left copy
    # Brand wordmark: "Context" ink + "Compiler" forest (matches .brand span)
    brand = slide.shapes.add_textbox(MARGIN, Inches(1.85), Inches(6.8), Inches(1.1))
    tf = brand.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Context "
    _set_run(r1, size=44, bold=True, color=INK, font=FONT_BRAND)
    r2 = p.add_run()
    r2.text = "Compiler"
    _set_run(r2, size=44, bold=True, color=FOREST, font=FONT_BRAND)
    _textbox(
        slide,
        MARGIN,
        Inches(2.85),
        Inches(6.8),
        Inches(0.35),
        [
            {
                "text": "MCP · LOCAL BM25 · NO KEY FOR COMPILE",
                "size": 11,
                "bold": False,
                "font": FONT_MONO,
                "color": FOREST,
            }
        ],
    )
    _textbox(
        slide,
        MARGIN,
        Inches(3.35),
        Inches(6.6),
        Inches(0.7),
        [
            {
                "text": "Same answer. 3% of the tokens.",
                "size": 26,
                "bold": True,
                "font": FONT_BODY,
                "color": INK,
            }
        ],
    )
    _textbox(
        slide,
        MARGIN,
        Inches(4.1),
        Inches(6.4),
        Inches(1.2),
        [
            {
                "text": "Task-aware compile under a hard token budget — coverage-first packing, omit honesty, MCP or the hosted web demo.",
                "size": 15,
                "color": MUTED,
            }
        ],
    )
    _textbox(
        slide,
        MARGIN,
        Inches(6.85),
        Inches(12),
        Inches(0.35),
        [
            {
                "text": "OpenAI × NamasteDev Codex Hackathon · July 2026 · solo build",
                "size": 11,
                "font": FONT_MONO,
                "color": FAINT,
            }
        ],
    )

    # Right: forest plane (savings visual)
    _rect(slide, Inches(8.05), Inches(1.7), Inches(4.55), Inches(4.0), PLANE)
    _textbox(
        slide,
        Inches(8.35),
        Inches(1.95),
        Inches(4.0),
        Inches(0.4),
        [
            {
                "text": "vendor_handbook.docx · one question",
                "size": 11,
                "font": FONT_MONO,
                "color": PLANE_MUTED,
            }
        ],
    )
    _textbox(
        slide,
        Inches(8.35),
        Inches(2.5),
        Inches(4.0),
        Inches(0.3),
        [{"text": "Whole file", "size": 13, "bold": True, "color": PLANE_INK}],
    )
    _textbox(
        slide,
        Inches(10.6),
        Inches(2.5),
        Inches(1.7),
        Inches(0.3),
        [
            {
                "text": "20,364 tokens",
                "size": 12,
                "font": FONT_MONO,
                "color": PLANE_INK,
                "align": PP_ALIGN.RIGHT,
            }
        ],
    )
    _rect(slide, Inches(8.35), Inches(2.85), Inches(3.95), Inches(0.14), RGBColor(0x2A, 0x4A, 0x3A))
    _rect(slide, Inches(8.35), Inches(2.85), Inches(3.95), Inches(0.14), WASTE)

    _textbox(
        slide,
        Inches(8.35),
        Inches(3.3),
        Inches(4.0),
        Inches(0.3),
        [{"text": "Compiled", "size": 13, "bold": True, "color": PLANE_INK}],
    )
    _textbox(
        slide,
        Inches(10.6),
        Inches(3.3),
        Inches(1.7),
        Inches(0.3),
        [
            {
                "text": "591 tokens",
                "size": 12,
                "font": FONT_MONO,
                "color": PLANE_INK,
                "align": PP_ALIGN.RIGHT,
            }
        ],
    )
    _rect(slide, Inches(8.35), Inches(3.65), Inches(3.95), Inches(0.14), RGBColor(0x2A, 0x4A, 0x3A))
    _rect(slide, Inches(8.35), Inches(3.65), Inches(0.16), Inches(0.14), COMPILED)  # ~3%

    _textbox(
        slide,
        Inches(8.35),
        Inches(4.15),
        Inches(4.0),
        Inches(0.9),
        [
            {
                "text": "97% fewer tokens",
                "size": 18,
                "bold": True,
                "color": PLANE_INK,
            },
            {
                "text": "Same facts. Every read.",
                "size": 13,
                "color": PLANE_MUTED,
                "space_before": 6,
            },
            {
                "text": "verified: unit tests + recall eval + live runs",
                "size": 10,
                "font": FONT_MONO,
                "color": PLANE_MUTED,
                "space_before": 10,
            },
        ],
    )


def slide_problem(prs: Presentation):
    slide = _blank_slide(prs)
    _textbox(
        slide,
        MARGIN,
        Inches(0.45),
        Inches(12),
        Inches(0.7),
        [
            {
                "text": "Agents pay to read pages they don’t need",
                "size": 28,
                "bold": True,
                "font": FONT_BRAND,
                "color": INK,
            }
        ],
    )

    rows = [
        (
            "Whole-file reads",
            "Agents convert and ingest entire documents to answer a single question — every time.",
            Inches(1.4),
        ),
        (
            "95%+ never touches the answer",
            "The relevant clause is 2 pages of 100. The other 98 are pure token spend.",
            Inches(2.9),
        ),
        (
            "And it compounds",
            "Every file × every question × every agent × every day. Context windows fill, latency and cost climb.",
            Inches(4.4),
        ),
    ]
    for title, body, top in rows:
        _rect(slide, MARGIN, top, Inches(0.08), Inches(1.1), FOREST)
        _textbox(
            slide,
            Inches(1.05),
            top,
            Inches(6.2),
            Inches(0.4),
            [{"text": title, "size": 18, "bold": True, "color": INK}],
        )
        _textbox(
            slide,
            Inches(1.05),
            top + Inches(0.4),
            Inches(6.2),
            Inches(0.7),
            [{"text": body, "size": 14, "color": MUTED}],
        )

    # Stat panels (hairline, surface — not neon cards)
    panel_h = Inches(2.15)
    for i, (big, small) in enumerate(
        [
            (
                "$61 → $2",
                "per 1,000 reads of one 100-page document (at $3/Mtok input — demo cost meter default)",
            ),
            (
                "95%",
                "of file tokens are wasted per task — across pdf, docx, xlsx, pptx, images",
            ),
        ]
    ):
        top = Inches(1.4) + i * (panel_h + Inches(0.25))
        _rect(slide, Inches(8.0), top, Inches(4.6), panel_h, SURFACE, line=BORDER)
        _accent_bar(slide, Inches(8.0), top, Inches(4.6))
        _textbox(
            slide,
            Inches(8.3),
            top + Inches(0.35),
            Inches(4.0),
            Inches(0.7),
            [{"text": big, "size": 32, "bold": True, "font": FONT_BRAND, "color": FOREST}],
        )
        _textbox(
            slide,
            Inches(8.3),
            top + Inches(1.15),
            Inches(4.0),
            Inches(0.8),
            [{"text": small, "size": 13, "color": MUTED}],
        )


def slide_market(prs: Presentation):
    slide = _blank_slide(prs)
    _textbox(
        slide,
        MARGIN,
        Inches(0.45),
        Inches(12),
        Inches(0.7),
        [
            {
                "text": "Who buys, why now, where we wedge",
                "size": 28,
                "bold": True,
                "font": FONT_BRAND,
                "color": INK,
            }
        ],
    )
    cards = [
        (
            Inches(0.7),
            Inches(1.4),
            "Who buys",
            "Teams wiring coding and document agents — Cursor, Claude Code, Codex, and internal agent platforms.",
        ),
        (
            Inches(6.85),
            Inches(1.4),
            "Why now",
            "MCP is the distribution rail. Context windows + $/Mtok make whole-file reads expensive. Agents re-read the same docs across tasks.",
        ),
        (
            Inches(0.7),
            Inches(4.15),
            "Wedge",
            "A task-budgeted file prep layer — not a full RAG platform. One job: select what the agent needs under a hard token budget.",
        ),
        (
            Inches(6.85),
            Inches(4.15),
            "Vs converters",
            "Converters: format-in → whole-file-out. We select under budget and return a recoverable omitted-sections manifest.",
        ),
    ]
    for left, top, title, body in cards:
        _rect(slide, left, top, Inches(5.8), Inches(2.4), SURFACE, line=BORDER)
        _accent_bar(slide, left, top, Inches(5.8))
        _textbox(
            slide,
            left + Inches(0.3),
            top + Inches(0.35),
            Inches(5.2),
            Inches(0.45),
            [{"text": title, "size": 18, "bold": True, "color": FOREST}],
        )
        _textbox(
            slide,
            left + Inches(0.3),
            top + Inches(0.9),
            Inches(5.2),
            Inches(1.3),
            [{"text": body, "size": 14, "color": MUTED}],
        )


def slide_how(prs: Presentation):
    slide = _blank_slide(prs)
    _textbox(
        slide,
        MARGIN,
        Inches(0.4),
        Inches(12),
        Inches(0.55),
        [
            {
                "text": "compile_context(file, task, budget)",
                "size": 26,
                "bold": True,
                "font": FONT_MONO,
                "color": INK,
            }
        ],
    )
    _textbox(
        slide,
        MARGIN,
        Inches(1.0),
        Inches(12),
        Inches(0.4),
        [
            {
                "text": "One tool call in — guaranteed-size, task-relevant markdown out.",
                "size": 15,
                "color": MUTED,
            }
        ],
    )

    steps = [
        ("1 · Convert", "markitdown + content-hash cache"),
        ("2 · Chunk", "heading-aware; tables never split"),
        ("3 · Rank", "BM25 + query cleanup — local, no LLM"),
        ("4 · Pack", "coverage-first under budget; content before manifest; relevance floor"),
    ]
    box_w = Inches(2.7)
    gap = Inches(0.22)
    start = MARGIN
    for i, (title, body) in enumerate(steps):
        left = start + i * (box_w + gap)
        _rect(slide, left, Inches(1.65), box_w, Inches(1.7), SURFACE, line=BORDER)
        _accent_bar(slide, left, Inches(1.65), box_w)
        _textbox(
            slide,
            left + Inches(0.18),
            Inches(1.9),
            box_w - Inches(0.3),
            Inches(0.4),
            [{"text": title, "size": 15, "bold": True, "color": FOREST, "font": FONT_MONO}],
        )
        _textbox(
            slide,
            left + Inches(0.18),
            Inches(2.4),
            box_w - Inches(0.3),
            Inches(0.8),
            [{"text": body, "size": 13, "color": MUTED}],
        )
        if i < 3:
            _rect(
                slide,
                left + box_w + Inches(0.02),
                Inches(2.4),
                Inches(0.18),
                Inches(0.04),
                BORDER2,
            )

    bottoms = [
        (
            Inches(0.7),
            "MCP + web demo",
            "Two tools: compile_context + expand_section. Cursor, Claude Code, Codex, Claude Desktop — plus the hosted demo. No API key for compile.",
        ),
        (
            Inches(6.85),
            "Omit honesty",
            "Every response ends with an omitted-sections manifest. Fetch any section by id via expand_section. Trimming is never silent.",
        ),
    ]
    for left, title, body in bottoms:
        _rect(slide, left, Inches(3.85), Inches(5.8), Inches(2.85), SURFACE, line=BORDER)
        _rect(slide, left, Inches(3.85), Inches(0.1), Inches(2.85), FOREST)
        _textbox(
            slide,
            left + Inches(0.4),
            Inches(4.15),
            Inches(5.1),
            Inches(0.45),
            [{"text": title, "size": 18, "bold": True, "color": INK}],
        )
        _textbox(
            slide,
            left + Inches(0.4),
            Inches(4.7),
            Inches(5.1),
            Inches(1.6),
            [{"text": body, "size": 14, "color": MUTED}],
        )


def slide_measured(prs: Presentation):
    slide = _blank_slide(prs)
    _textbox(
        slide,
        MARGIN,
        Inches(0.45),
        Inches(12),
        Inches(0.7),
        [
            {
                "text": "Measured, not promised",
                "size": 28,
                "bold": True,
                "font": FONT_BRAND,
                "color": INK,
            }
        ],
    )
    stats = [
        ("97.1%", "token reduction on a real 78-section docx, answer intact"),
        ("34×", "cheaper per read — cached, so repeat reads are instant"),
        ("Parity", "Full file vs compile (+ Include in Prove) — dual buttons; not an Agent run"),
    ]
    w = Inches(3.85)
    for i, (big, small) in enumerate(stats):
        left = MARGIN + i * (w + Inches(0.2))
        _rect(slide, left, Inches(1.4), w, Inches(2.35), SURFACE, line=BORDER)
        _accent_bar(slide, left, Inches(1.4), w)
        _textbox(
            slide,
            left + Inches(0.25),
            Inches(1.75),
            w - Inches(0.4),
            Inches(0.7),
            [{"text": big, "size": 32, "bold": True, "font": FONT_BRAND, "color": FOREST}],
        )
        _textbox(
            slide,
            left + Inches(0.25),
            Inches(2.55),
            w - Inches(0.4),
            Inches(0.95),
            [{"text": small, "size": 13, "color": MUTED}],
        )

    # Example bar strip on forest plane
    _rect(slide, MARGIN, Inches(4.15), Inches(11.9), Inches(2.55), PLANE)
    _textbox(
        slide,
        Inches(1.0),
        Inches(4.4),
        Inches(11.2),
        Inches(0.4),
        [
            {
                "text": 'vendor_handbook.docx · “How long do refunds take and who approves large ones?” · budget 1,200',
                "size": 12,
                "font": FONT_MONO,
                "color": PLANE_MUTED,
            }
        ],
    )
    _textbox(
        slide,
        Inches(1.0),
        Inches(5.0),
        Inches(11.2),
        Inches(0.3),
        [
            {
                "text": "Whole file — 20,364 tokens",
                "size": 13,
                "bold": True,
                "color": PLANE_INK,
            }
        ],
    )
    _rect(slide, Inches(1.0), Inches(5.35), Inches(10.8), Inches(0.16), RGBColor(0x2A, 0x4A, 0x3A))
    _rect(slide, Inches(1.0), Inches(5.35), Inches(10.8), Inches(0.16), WASTE)
    _textbox(
        slide,
        Inches(1.0),
        Inches(5.7),
        Inches(11.2),
        Inches(0.3),
        [
            {
                "text": "Compiled — 591 tokens · answer intact · verified by test suite + recall eval + MCP run",
                "size": 13,
                "bold": True,
                "color": PLANE_INK,
            }
        ],
    )
    _rect(slide, Inches(1.0), Inches(6.05), Inches(10.8), Inches(0.16), RGBColor(0x2A, 0x4A, 0x3A))
    _rect(slide, Inches(1.0), Inches(6.05), Inches(0.35), Inches(0.16), COMPILED)


def slide_next(prs: Presentation):
    slide = _blank_slide(prs)
    _textbox(
        slide,
        MARGIN,
        Inches(0.45),
        Inches(12),
        Inches(0.7),
        [
            {
                "text": "What’s next",
                "size": 28,
                "bold": True,
                "font": FONT_BRAND,
                "color": INK,
            }
        ],
    )
    items = [
        (
            "Local embeddings",
            "Query cleanup + offline recall eval shipped; local embeddings next for paraphrase.",
        ),
        (
            "Multi-file corpora",
            "Compile across folders; shared conversion-cache gateway for teams.",
        ),
        (
            "Video & audio",
            "Transcription into the same pipeline (cut for scope, not feasibility).",
        ),
        (
            "npx context-compiler",
            "One-command install for any MCP client.",
        ),
    ]
    for i, (title, body) in enumerate(items):
        top = Inches(1.35) + i * Inches(1.15)
        _rect(slide, MARGIN, top, Inches(0.08), Inches(0.9), FOREST)
        _textbox(
            slide,
            Inches(1.05),
            top,
            Inches(6.3),
            Inches(0.35),
            [{"text": title, "size": 16, "bold": True, "color": INK}],
        )
        _textbox(
            slide,
            Inches(1.05),
            top + Inches(0.35),
            Inches(6.3),
            Inches(0.55),
            [{"text": body, "size": 13, "color": MUTED}],
        )

    # CTA plane
    _rect(slide, Inches(8.0), Inches(1.35), Inches(4.6), Inches(5.2), PLANE)
    _textbox(
        slide,
        Inches(8.35),
        Inches(1.75),
        Inches(4.0),
        Inches(0.45),
        [{"text": "Try it live", "size": 22, "bold": True, "font": FONT_BRAND, "color": PLANE_INK}],
    )
    _textbox(
        slide,
        Inches(8.35),
        Inches(2.4),
        Inches(4.0),
        Inches(1.2),
        [
            {
                "text": "demo:  context-compiler.onrender.com\n       Docker / Render deploy of this repo\nrepo:  github.com/shenba1712/context-compiler",
                "size": 12,
                "font": FONT_MONO,
                "color": PLANE_MUTED,
            }
        ],
    )
    _textbox(
        slide,
        Inches(8.35),
        Inches(3.85),
        Inches(4.0),
        Inches(1.2),
        [
            {
                "text": "Live demo + MCP in Cursor / Claude Code / Codex. Same two tools; prove answer parity on a real doc.",
                "size": 14,
                "color": PLANE_INK,
            }
        ],
    )
    _textbox(
        slide,
        Inches(8.35),
        Inches(5.4),
        Inches(4.0),
        Inches(0.8),
        [
            {
                "text": "Stop paying for pages your agent doesn’t read.",
                "size": 16,
                "bold": True,
                "color": COMPILED,
            }
        ],
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Remove default empty slide if any
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    slide_title(prs)
    slide_problem(prs)
    slide_market(prs)
    slide_how(prs)
    slide_measured(prs)
    slide_next(prs)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
